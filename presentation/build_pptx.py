"""
AUC Data Challenge 2026 - Presentation Slide Deck (PPTX)
Team DataRx | St. Helena Parish, Louisiana | Census Tract 22091 951100
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(BASE, 'outputs')
OUT = os.path.join(BASE, 'presentation')
os.makedirs(OUT, exist_ok=True)

NAVY = RGBColor(29, 53, 87)
TEAL = RGBColor(42, 157, 143)
CORAL = RGBColor(230, 57, 70)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(141, 153, 174)
DARK = RGBColor(51, 51, 51)
LIGHT_BG = RGBColor(248, 249, 250)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # blank layout


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color=TEAL):
    shp = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullets(slide, items, x, y, w, h, size=14, color=DARK, bold_prefix=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.level = 0

        if bold_prefix and ':' in item:
            prefix, rest = item.split(':', 1)
            run1 = p.add_run()
            run1.text = prefix + ':'
            run1.font.size = Pt(size)
            run1.font.bold = True
            run1.font.color.rgb = TEAL
            run1.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
    return tf


def add_image(slide, path, x, y, w=None, h=None):
    full = os.path.join(IMG, path) if not os.path.isabs(path) else path
    if os.path.exists(full):
        kwargs = {}
        if w: kwargs['width'] = Inches(w)
        if h: kwargs['height'] = Inches(h)
        slide.shapes.add_picture(full, Inches(x), Inches(y), **kwargs)
    else:
        add_text(slide, f'[Missing: {path}]', x, y, 3, 0.5, size=10, color=CORAL)


def slide_number(slide, n, total=11):
    add_text(slide, f'{n}/{total}', 12.5, 7.1, 0.8, 0.3, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def source_text(slide, text):
    add_text(slide, text, 0.5, 7.0, 12, 0.3, size=7, color=MUTED)


# ======================================================================
# SLIDE 1: Title
# ======================================================================
print("Slide 1: Title")
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_accent_bar(s, TEAL)

add_text(s, 'Community Hub Marketplace', 1, 1.5, 11, 1, size=44, bold=True, color=WHITE)
add_text(s, 'Bridging Economic Gaps to Improve Healthcare Access', 1, 2.5, 11, 0.6, size=22, color=TEAL)
add_text(s, 'St. Helena Parish, Louisiana  |  Census Tract 22091 951100', 1, 3.6, 11, 0.5, size=16, color=RGBColor(180, 180, 195))
add_text(s, 'AUC Data Challenge 2026', 1, 4.1, 11, 0.5, size=16, color=RGBColor(180, 180, 195))
add_text(s, 'Team DataRx', 1, 5.2, 11, 0.6, size=24, bold=True, color=TEAL)
add_text(s, '"In communities with an IGS below 45, how do economic conditions affect healthcare access, and how can small businesses help improve outcomes?"',
         1, 6.2, 11, 0.8, size=11, color=RGBColor(180, 180, 195))
slide_number(s, 1)

# ======================================================================
# SLIDE 2: Objective & Problem Statement
# ======================================================================
print("Slide 2: Objective & Problem Statement")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Objective & Problem Statement', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)

add_text(s, 'The Challenge:', 0.8, 1.3, 5, 0.4, size=18, bold=True, color=DARK)
add_text(s, 'St. Helena Parish has an Inclusive Growth Score of just 36 -- well below the 45 threshold. '
         'This rural community of ~10,000 faces compounding economic and health crises.',
         0.8, 1.8, 5.5, 1.0, size=13, color=DARK)

add_bullets(s, [
    'No local jobs: 62% of working-age adults not employed',
    'Worker exodus: 90% commute out of parish for work',
    'No insurance: only 7.9% have employer coverage (vs 55% nationally)',
    'Healthcare desert: zero pharmacies, 1 dentist, 2 PCPs for 10,000 people',
    'No internet: 3rd percentile nationally for internet access',
    'Deep poverty: 33.7% poverty rate (2.7x the national rate)',
], 0.8, 3.0, 5.5, 2.5, size=12, color=DARK)

add_text(s, 'Our Solution: Community Hub Marketplace', 0.8, 5.6, 6, 0.4, size=18, bold=True, color=TEAL)
add_text(s, 'A physical and digital center that creates local jobs with employer insurance, '
         'brings healthcare services, provides broadband, incubates small businesses, '
         'and keeps spending local. One intervention, multiple IGS levers.',
         0.8, 6.1, 5.5, 1.0, size=12, color=DARK)

add_image(s, 'po/po_hub_metrics.png', 7.2, 1.2, w=5.5)
slide_number(s, 2)

# ======================================================================
# SLIDE 3: Community Overview
# ======================================================================
print("Slide 3: Community Overview")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Community Overview', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)
add_text(s, 'Demographics, Economic Conditions & Social Metrics', 0.8, 0.9, 11, 0.4, size=14, color=MUTED)

add_image(s, 'po/po_economic_table.png', 0.3, 1.5, w=6)
add_image(s, 'po/po_igs_breakdown.png', 7, 1.5, w=5.5)

add_image(s, 'e1/e1_poverty_rate.png', 0.3, 4.0, w=4)
add_image(s, 'e1/e1_income_comparison.png', 4.5, 4.0, w=4)
add_image(s, 'ag/ag_food_paradox.png', 8.7, 4.0, w=4)

source_text(s, 'Source: U.S. Census ACS 2023, Inclusive Growth Score 2025  |  api.census.gov  |  inclusivegrowthscore.com')
slide_number(s, 3)

# ======================================================================
# SLIDE 4: Benchmarking
# ======================================================================
print("Slide 4: Benchmarking")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Benchmarking', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)
add_text(s, 'IGS Ranking vs State and National Averages', 0.8, 0.9, 11, 0.4, size=14, color=MUTED)

add_image(s, 'po/po_igs_history.png', 0.3, 1.5, w=5.5)
add_image(s, 'po/po_igs_breakdown.png', 6.5, 1.5, w=5.5)
add_image(s, 'bm/bm_benchmarking.png', 1, 4.2, w=11)

source_text(s, 'Source: Inclusive Growth Score 2017-2025  |  inclusivegrowthscore.com')
slide_number(s, 4)

# ======================================================================
# SLIDE 5: Key Findings - Economic
# ======================================================================
print("Slide 5: Key Findings - Economic")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Key Findings: Economic Crisis', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)
add_text(s, 'Poverty, Employment & the Worker Exodus', 0.8, 0.9, 11, 0.4, size=14, color=MUTED)

add_image(s, 'e1/e1_poverty_by_age.png', 0.2, 1.4, w=4.3)
add_image(s, 'e2/e2_labor_force.png', 4.6, 1.4, w=4.3)
add_image(s, 'e2/e2_worker_exodus.png', 9.0, 1.4, w=4.3)

add_image(s, 'e2/e2_commute_times.png', 0.2, 4.2, w=4.3)
add_image(s, 'e3/e3_insurance_by_employment.png', 4.6, 4.2, w=4.3)
add_image(s, 'e3/e3_industry_by_gender.png', 9.0, 4.2, w=4.3)

slide_number(s, 5)

# ======================================================================
# SLIDE 6: Key Findings - Healthcare
# ======================================================================
print("Slide 6: Key Findings - Healthcare")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Key Findings: Healthcare Desert', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)
add_text(s, 'Disease Burden, Provider Shortages & Preventable Hospitalizations', 0.8, 0.9, 11, 0.4, size=14, color=MUTED)

add_image(s, 'h1/h1_chronic_disease.png', 0.2, 1.4, w=4.3)
add_image(s, 'h2/h2_provider_counts.png', 4.6, 1.4, w=4.3)
add_image(s, 'h2/h2_preventable_hosp.png', 9.0, 1.4, w=4.3)

add_image(s, 'h1/h1_social_determinants.png', 0.2, 4.2, w=4.3)
add_image(s, 'h2/h2_life_expectancy.png', 4.6, 4.2, w=4.3)
add_image(s, 's1/s1_healthcare_gaps.png', 9.0, 4.2, w=4.3)

slide_number(s, 6)

# ======================================================================
# SLIDE 7: Proposed Solution
# ======================================================================
print("Slide 7: Proposed Solution")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Proposed Solution: Community Hub Marketplace', 0.8, 0.3, 11, 0.8, size=30, bold=True, color=NAVY)
add_text(s, 'One Hub. Multiple IGS Levers.', 0.8, 1.1, 6, 0.4, size=18, bold=True, color=TEAL)

add_bullets(s, [
    'Local Jobs: Create 25-50 jobs with employer-sponsored health insurance',
    'Broadband Center: Free WiFi + computer lab (internet: 3rd to 23rd percentile)',
    'Pharmacy & Clinic: First pharmacy in parish; telehealth clinic referrals',
    'Small Business Incubator: Vendor stalls, shared kitchen, SBA loan help',
    'Fresh Food Market: Local produce combating 33.4% food insecurity',
    'Daycare & After-School: Enabling workforce participation for parents',
    'Digital Services: Job portal, government kiosk, financial literacy',
], 0.8, 1.6, 5.5, 3.5, size=12, color=DARK, bold_prefix=True)

add_text(s, 'The Hub targets 12 of 18 IGS drivers', 0.8, 5.4, 5, 0.4, size=14, bold=True, color=NAVY)

add_image(s, 'po/po_feature_importance.png', 7, 1.0, w=6)

source_text(s, 'Source: XGBoost model feature importance trained on 1.4M IGS observations')
slide_number(s, 7)

# ======================================================================
# SLIDE 8: Implementation Plan
# ======================================================================
print("Slide 8: Implementation Plan")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Implementation Plan', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)

# Year 1
add_text(s, 'Year 1: Foundation', 0.8, 1.2, 4, 0.4, size=16, bold=True, color=TEAL)
add_bullets(s, [
    'Secure facility (repurpose vacant commercial building)',
    'Launch broadband center + digital services kiosk',
    'Open fresh food market with local produce',
    'Partner with FQHCs for telehealth clinic',
    'Begin SBA loan outreach for vendors',
], 0.8, 1.7, 3.8, 2.5, size=11, color=DARK)

# Year 2-3
add_text(s, 'Year 2-3: Growth', 4.8, 1.2, 4, 0.4, size=16, bold=True, color=TEAL)
add_bullets(s, [
    'Recruit pharmacy operator',
    'Expand to 15-20 vendor stalls',
    'Launch daycare / after-school program',
    'Employer insurance for all Hub employees',
    'Digital literacy & job training programs',
], 4.8, 1.7, 3.8, 2.5, size=11, color=DARK)

# Year 4-5
add_text(s, 'Year 4-5: Scale', 8.8, 1.2, 4, 0.4, size=16, bold=True, color=TEAL)
add_bullets(s, [
    'Hub reaches 40-50 employees with benefits',
    'Second location or mobile Hub',
    'Mentor businesses to independent storefronts',
    'Community health worker program',
    'Target: IGS above 45 threshold',
], 8.8, 1.7, 3.8, 2.5, size=11, color=DARK)

# Stakeholders
add_text(s, 'Key Stakeholders', 0.8, 4.5, 12, 0.4, size=16, bold=True, color=NAVY)
add_text(s, 'St. Helena Parish Government  |  SBA  |  USDA Rural Development  |  FQHCs  |  Mastercard Center for Inclusive Growth',
         0.8, 5.0, 11, 0.4, size=12, color=DARK)

# Risks
add_text(s, 'Risks & Mitigations', 0.8, 5.6, 12, 0.4, size=16, bold=True, color=NAVY)
add_bullets(s, [
    'Low foot traffic: mitigate with anchor tenants + essential services',
    'Funding gaps: mitigate with USDA grants + New Market Tax Credits',
    'Workforce readiness: mitigate with on-site training programs',
], 0.8, 6.1, 11, 1.2, size=11, color=DARK)

slide_number(s, 8)

# ======================================================================
# SLIDE 9: Predicted Outcomes
# ======================================================================
print("Slide 9: Predicted Outcomes")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Predicted Outcomes', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)
add_text(s, 'XGBoost Model (R-squared=0.99) trained on 1.4M IGS observations across 85K tracts', 0.8, 0.9, 11, 0.4, size=14, color=MUTED)

add_image(s, 'po/po_trajectory.png', 0.2, 1.4, w=6.5)
add_image(s, 'po/po_radar.png', 6.8, 1.3, w=6)

add_image(s, 'po/po_scenario_bars.png', 0.2, 4.2, w=6.5)
add_image(s, 'po/po_waterfall.png', 6.8, 4.2, w=6.2)

slide_number(s, 9)

# ======================================================================
# SLIDE 10: Metrics for Success
# ======================================================================
print("Slide 10: Metrics for Success")
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_accent_bar(s)

add_text(s, 'Metrics for Success', 0.8, 0.3, 11, 0.8, size=32, bold=True, color=NAVY)
add_text(s, 'Indicators to Track Long-Term Impact', 0.8, 0.9, 11, 0.4, size=14, color=MUTED)

add_bullets(s, [
    'IGS Score: 36 to 45+ (primary outcome, model-validated)',
    'Employer Insurance Rate: 7.9% to 25% (Hub jobs with benefits)',
    'Labor Participation: 50.9% to 60% (local job creation)',
    'Internet Access: 3rd to 25th percentile (broadband center)',
    'Commercial Diversity: 34th to 54th percentile (new business sectors)',
    'Preventable Hospitalizations: 4,862 to 3,500 per 100K (local healthcare)',
    'Food Insecurity: 33.4% to 20% (fresh food market)',
], 0.8, 1.5, 5.5, 3.5, size=13, color=DARK, bold_prefix=True)

add_image(s, 'ms/ms_metrics_dashboard.png', 7, 1.3, w=6)

source_text(s, 'Source: XGBoost simulation + Census ACS + CDC PLACES + County Health Rankings')
slide_number(s, 10)

# ======================================================================
# SLIDE 11: Conclusion
# ======================================================================
print("Slide 11: Conclusion")
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_accent_bar(s, TEAL)

add_text(s, 'Conclusion', 1, 1.0, 11, 1, size=38, bold=True, color=WHITE)
add_text(s, 'St. Helena Parish does not need charity -- it needs infrastructure.\nThe Community Hub Marketplace is that infrastructure.',
         1, 2.0, 11, 1, size=18, color=TEAL)

add_bullets(s, [
    'One intervention targeting 12 of 18 IGS drivers simultaneously',
    'XGBoost model (R-squared=0.99) projects IGS from 36 to 46.5 in 5 years',
    'Crossing the 45 threshold -- answering the challenge question directly',
    'Employer insurance alone could 3x coverage (7.9% to 25%)',
    '10 real Louisiana tracts already at our Year 3 target -- this is achievable',
    'Every dollar spent locally recirculates -- the Hub creates a virtuous cycle',
], 1, 3.5, 11, 3, size=15, color=RGBColor(200, 200, 210))

add_text(s, 'Team DataRx', 1, 6.0, 11, 0.5, size=22, bold=True, color=TEAL)
add_text(s, 'AUC Data Challenge 2026  |  St. Helena Parish, Louisiana', 1, 6.5, 11, 0.4, size=12, color=RGBColor(180, 180, 195))

slide_number(s, 11, 11)

# ─── Save ─────────────────────────────────────────────────────────────────────
outpath = os.path.join(OUT, 'DataRx_AUC_Challenge_2026.pptx')
prs.save(outpath)
print(f'\nDone! Saved: {outpath}')
print(f'Slides: 11 (Title + 10 content)')
